"""
rKYC 해커톤 발표 자료 PPT 생성 스크립트
python-pptx를 사용하여 PowerPoint 파일 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 색상 정의
DARK_BG = RGBColor(17, 24, 39)  # 진한 남색
ACCENT_BLUE = RGBColor(59, 130, 246)  # 밝은 파랑
ACCENT_PURPLE = RGBColor(139, 92, 246)  # 보라
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(156, 163, 175)
SUCCESS_GREEN = RGBColor(34, 197, 94)
WARNING_YELLOW = RGBColor(250, 204, 21)
DANGER_RED = RGBColor(239, 68, 68)
DARK_CARD = RGBColor(31, 41, 55)

def set_slide_background(slide, color):
    """슬라이드 배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """타이틀 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, highlight_indices=None):
    """콘텐츠 슬라이드 추가 (제목 + 불릿)"""
    if highlight_indices is None:
        highlight_indices = []

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # 불릿 포인트
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"  {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = ACCENT_BLUE if i in highlight_indices else LIGHT_GRAY
        p.space_before = Pt(12)
        p.level = 0

    return slide

def add_architecture_slide(prs):
    """아키텍처 다이어그램 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 박스들 그리기
    boxes = [
        ("Frontend\n(React/Vercel)", 0.5, 1.3, 2.5, 1.0, ACCENT_BLUE),
        ("Backend\n(FastAPI/Railway)", 3.5, 1.3, 2.5, 1.0, ACCENT_PURPLE),
        ("Worker\n(Celery/Redis)", 6.5, 1.3, 2.5, 1.0, SUCCESS_GREEN),
        ("Supabase\nPostgreSQL", 2.0, 3.0, 2.5, 1.0, WARNING_YELLOW),
        ("LLM APIs\n(4 Models)", 5.5, 3.0, 2.5, 1.0, DANGER_RED),
    ]

    for text, x, y, w, h, color in boxes:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 설명
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    tf = desc_box.text_frame
    tf.word_wrap = True

    descs = [
        "Frontend: API 호출만 (LLM 직접 호출 X)",
        "Backend: DB CRUD만 (LLM 직접 호출 X)",
        "Worker: 모든 LLM 호출 담당 (보안)",
        "핵심: LLM 키 격리로 보안 강화"
    ]

    for i, desc in enumerate(descs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)

    return slide

def add_multi_agent_slide(prs):
    """Multi-Agent 아키텍처 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Agent Consensus Engine"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 에이전트 박스들
    agents = [
        ("Perplexity\nLayer 1: Search", 0.5, 1.5, 2.8, 1.2, ACCENT_BLUE),
        ("Gemini\nLayer 1.5: Validate", 3.6, 1.5, 2.8, 1.2, ACCENT_PURPLE),
        ("Claude Opus\nLayer 2: Synthesize", 6.7, 1.5, 2.8, 1.2, SUCCESS_GREEN),
    ]

    for text, x, y, w, h, color in agents:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Consensus 박스
    consensus = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(3.2), Inches(6), Inches(1.5))
    consensus.fill.solid()
    consensus.fill.fore_color.rgb = WARNING_YELLOW
    consensus.line.fill.background()

    tf = consensus.text_frame
    p = tf.paragraphs[0]
    p.text = "Consensus Engine\nHybrid Similarity (Embedding + Jaccard)\n한국어 형태소 분석 (kiwipiepy)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER

    # 설명
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1.5))
    tf = desc_box.text_frame
    tf.word_wrap = True

    descs = [
        "일치 (유사도 >= 70%): Perplexity 값 채택, HIGH 신뢰도",
        "불일치: Perplexity 값 채택 + discrepancy 플래그 표시"
    ]

    for i, desc in enumerate(descs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(8)

    return slide

def add_fallback_slide(prs):
    """4-Layer Fallback 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "4-Layer Fallback: Never Fail"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 레이어 박스들
    layers = [
        ("Layer 0: Cache", "7일 TTL, 기존 프로필 재사용", 0.5, 1.3, SUCCESS_GREEN),
        ("Layer 1+1.5: Perplexity + Gemini", "실시간 검색 + 교차 검증", 0.5, 2.3, ACCENT_BLUE),
        ("Layer 2: Claude Synthesis", "Consensus Engine 합성", 0.5, 3.3, ACCENT_PURPLE),
        ("Layer 3: Rule-Based Merge", "LLM 없이 규칙 기반 병합", 0.5, 4.3, WARNING_YELLOW),
        ("Layer 4: Graceful Degradation", "최소 프로필 + 경고 플래그", 0.5, 5.3, DANGER_RED),
    ]

    for title_text, desc, x, y, color in layers:
        # 박스
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(9), Inches(0.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{title_text}  |  {desc}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE if color != WARNING_YELLOW else DARK_BG
        p.alignment = PP_ALIGN.LEFT

    return slide

def add_anti_hallucination_slide(prs):
    """Anti-Hallucination 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Anti-Hallucination: 4-Layer Defense"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 레이어 설명
    layers = [
        ("1. Source Verification", "Perplexity Citation 검증, 도메인 신뢰도 분류"),
        ("2. Extraction Guardrails", "LLM 프롬프트: \"모르면 null 반환\" 규칙"),
        ("3. Validation Layer", "Gemini 교차 검증, 범위/일관성 체크"),
        ("4. Audit Trail", "필드별 source_url, excerpt, confidence 추적"),
    ]

    y_pos = 1.3
    for title_text, desc in layers:
        # 제목 박스
        title_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(4), Inches(0.7))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = ACCENT_BLUE
        title_shape.line.fill.background()

        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # 설명 박스
        desc_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(y_pos), Inches(4.8), Inches(0.7))
        desc_shape.fill.solid()
        desc_shape.fill.fore_color.rgb = DARK_CARD
        desc_shape.line.fill.background()

        tf = desc_shape.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

        y_pos += 1.1

    # 결과 강조
    result_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.8))
    tf = result_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Result: 모든 시그널에 Evidence(출처) 필수"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_pipeline_slide(prs):
    """9-Stage Pipeline 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "9-Stage Analysis Pipeline"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 파이프라인 단계
    stages = [
        ("1", "SNAPSHOT", SUCCESS_GREEN),
        ("2", "DOC", ACCENT_BLUE),
        ("3", "PROFILE", ACCENT_PURPLE),
        ("4", "EXTERNAL", WARNING_YELLOW),
        ("5", "CONTEXT", LIGHT_GRAY),
        ("6", "SIGNAL", DANGER_RED),
        ("7", "VALIDATE", SUCCESS_GREEN),
        ("8", "INDEX", ACCENT_BLUE),
        ("9", "INSIGHT", ACCENT_PURPLE),
    ]

    x_pos = 0.3
    for num, name, color in stages:
        # 단계 박스
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.5), Inches(1.0), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE if color not in [WARNING_YELLOW, LIGHT_GRAY] else DARK_BG
        p.alignment = PP_ALIGN.CENTER

        # 라벨
        label_box = slide.shapes.add_textbox(Inches(x_pos - 0.1), Inches(2.6), Inches(1.2), Inches(0.8))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        x_pos += 1.05

    # AI 모델 매핑 표
    table_data = [
        ("Stage", "AI Model", "역할"),
        ("DOC_INGEST", "Claude Vision", "PDF/이미지 분석"),
        ("PROFILING", "Perplexity+Gemini", "기업 프로필 수집"),
        ("EXTERNAL", "Perplexity", "뉴스/이벤트 검색"),
        ("SIGNAL", "Claude Opus", "시그널 추출"),
        ("INSIGHT", "Claude Opus", "인사이트 생성"),
    ]

    y_pos = 3.5
    for i, (col1, col2, col3) in enumerate(table_data):
        color = WHITE if i == 0 else LIGHT_GRAY
        bold = i == 0

        for j, text in enumerate([col1, col2, col3]):
            text_box = slide.shapes.add_textbox(Inches(0.5 + j * 3), Inches(y_pos), Inches(2.8), Inches(0.4))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(14)
            p.font.bold = bold
            p.font.color.rgb = color

        y_pos += 0.4

    return slide

def add_circuit_breaker_slide(prs):
    """Circuit Breaker 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Circuit Breaker Pattern"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 상태 박스
    states = [
        ("CLOSED", "정상 동작\n모든 요청 허용", SUCCESS_GREEN, 0.5),
        ("OPEN", "차단됨\n5분 대기", DANGER_RED, 3.5),
        ("HALF_OPEN", "테스트 중\n1개 요청만", WARNING_YELLOW, 6.5),
    ]

    for name, desc, color, x in states:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.5), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{name}\n\n{desc}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE if color != WARNING_YELLOW else DARK_BG
        p.alignment = PP_ALIGN.CENTER

    # Provider 설정 표
    table_data = [
        ("Provider", "Threshold", "Cooldown"),
        ("Perplexity", "3회", "5분"),
        ("Gemini", "3회", "5분"),
        ("Claude", "2회", "10분"),
    ]

    y_pos = 3.5
    for i, (col1, col2, col3) in enumerate(table_data):
        color = WHITE if i == 0 else LIGHT_GRAY
        bold = i == 0

        for j, text in enumerate([col1, col2, col3]):
            text_box = slide.shapes.add_textbox(Inches(2 + j * 2), Inches(y_pos), Inches(1.8), Inches(0.4))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(16)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER

        y_pos += 0.45

    # Redis 설명
    redis_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(0.5))
    tf = redis_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Redis 영속화: Worker 재시작 시에도 상태 유지"
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_tech_stack_slide(prs):
    """기술 스택 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Technology Stack"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 스택 항목
    stacks = [
        ("Frontend", "React 18, TypeScript, TanStack Query, shadcn/ui", ACCENT_BLUE),
        ("Backend", "FastAPI, SQLAlchemy 2.0, Pydantic v2", ACCENT_PURPLE),
        ("Worker", "Celery, Redis, litellm (Multi-provider)", SUCCESS_GREEN),
        ("Database", "Supabase PostgreSQL, pgvector (HNSW)", WARNING_YELLOW),
        ("AI Models", "Claude Opus, GPT-5.2 Pro, Gemini 3 Pro, Perplexity", DANGER_RED),
        ("Deploy", "Vercel (Frontend), Railway (Backend/Worker)", LIGHT_GRAY),
    ]

    y_pos = 1.2
    for name, tech, color in stacks:
        # 카테고리
        cat_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(2), Inches(0.65))
        cat_shape.fill.solid()
        cat_shape.fill.fore_color.rgb = color
        cat_shape.line.fill.background()

        tf = cat_shape.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE if color not in [WARNING_YELLOW, LIGHT_GRAY] else DARK_BG
        p.alignment = PP_ALIGN.CENTER

        # 기술
        tech_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(6.8), Inches(0.65))
        tf = tech_box.text_frame
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        tf.paragraphs[0].space_before = Pt(8)

        y_pos += 0.8

    return slide

def add_impact_slide(prs):
    """비즈니스 임팩트 슬라이드 (ROI 강화)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Business Impact & ROI"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ROI 계산 박스
    roi_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(1.2))
    roi_box.fill.solid()
    roi_box.fill.fore_color.rgb = DARK_CARD
    roi_box.line.color.rgb = SUCCESS_GREEN
    roi_box.line.width = Pt(2)

    roi_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.9))
    tf = roi_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💰 시간 절감 계산"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p = tf.add_paragraph()
    p.text = "심사역 100명 × 월 10건 × 2시간 절감 = 월 2,000시간 = 연간 24,000시간"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Before/After 비교표
    table_y = 2.5
    headers = [("항목", 0.5), ("Before", 3.0), ("After", 5.5), ("효과", 7.5)]
    for text, x in headers:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(table_y), Inches(2.3), Inches(0.5))
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_BLUE
        box.line.fill.background()
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    rows = [
        ("외부정보 수집", "기업당 2시간", "10분 (자동)", "92% 절감"),
        ("리스크 모니터링", "수동/주기적", "실시간 알림", "24/7 상시"),
        ("KYC 갱신 판단", "담당자 경험", "시그널 트리거", "표준화"),
    ]

    for i, row in enumerate(rows):
        row_y = table_y + 0.5 + (i * 0.5)
        colors = [WHITE, LIGHT_GRAY, SUCCESS_GREEN, ACCENT_BLUE]
        for j, (text, x) in enumerate(zip(row, [0.5, 3.0, 5.5, 7.5])):
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(row_y), Inches(2.3), Inches(0.5))
            box.fill.solid()
            box.fill.fore_color.rgb = DARK_CARD
            box.line.fill.background()
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(12)
            p.font.color.rgb = colors[j]
            p.alignment = PP_ALIGN.CENTER

    # 정성적 효과
    qual_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(9), Inches(1.8))
    tf = qual_box.text_frame
    tf.word_wrap = True

    quals = [
        "✓ 담당자 업무 부담 감소 → 전략적 업무에 집중",
        "✓ 일관된 분석 품질 → 담당자별 편차 해소",
        "✓ 기회 시그널도 포착 → 리스크만이 아닌 균형 분석"
    ]

    for i, qual in enumerate(quals):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = qual
        p.font.size = Pt(16)
        p.font.color.rgb = SUCCESS_GREEN
        p.space_before = Pt(6)

    return slide


def add_competitive_slide(prs):
    """경쟁 제출작 대비 차별점 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "vs 일반적인 해커톤 제출작"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 비교표
    comparisons = [
        ("일반 제출작", "rKYC", DANGER_RED, SUCCESS_GREEN),
    ]

    # 헤더
    header_y = 1.1
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(header_y), Inches(4.3), Inches(0.6))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = DANGER_RED
    left_box.line.fill.background()
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ 일반적인 해커톤 제출작"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(header_y), Inches(4.3), Inches(0.6))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = SUCCESS_GREEN
    right_box.line.fill.background()
    tf = right_box.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ rKYC"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 비교 항목들
    items = [
        ('"GPT에게 물어보기" 챗봇', "9-stage 구조화된 파이프라인"),
        ("단일 LLM 의존", "Multi-Agent + 4-Layer Fallback"),
        ("프로토타입 수준", "실제 배포된 서비스 (Vercel+Railway)"),
        ("아이디어 발표", "작동하는 데모"),
        ("LLM 호출 = 끝", "Anti-Hallucination 4중 방어"),
        ("장애 시 서비스 중단", "Circuit Breaker로 무중단"),
    ]

    y_pos = 1.85
    for left_text, right_text in items:
        # 왼쪽 (일반)
        left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(4.3), Inches(0.55))
        left.fill.solid()
        left.fill.fore_color.rgb = DARK_CARD
        left.line.fill.background()
        tf = left.text_frame
        p = tf.paragraphs[0]
        p.text = left_text
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

        # 오른쪽 (rKYC)
        right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(y_pos), Inches(4.3), Inches(0.55))
        right.fill.solid()
        right.fill.fore_color.rgb = DARK_CARD
        right.line.color.rgb = SUCCESS_GREEN
        right.line.width = Pt(1)
        tf = right.text_frame
        p = tf.paragraphs[0]
        p.text = right_text
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_GREEN
        p.alignment = PP_ALIGN.CENTER

        y_pos += 0.65

    # 완성도 강조
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.9), Inches(6), Inches(0.7))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = ACCENT_BLUE
    highlight.line.fill.background()
    tf = highlight.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 해커톤에서 이 정도 완성도는 상위 5%"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_integration_slide(prs):
    """내부 시스템 연계 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "기존 인프라와의 연계"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 여신시스템 박스
    loan_sys = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(2.5), Inches(1.2))
    loan_sys.fill.solid()
    loan_sys.fill.fore_color.rgb = ACCENT_BLUE
    loan_sys.line.fill.background()
    tf = loan_sys.text_frame
    p = tf.paragraphs[0]
    p.text = "여신시스템\n\n고객번호, 여신잔액\n담보정보, 내부등급"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # KYC시스템 박스
    kyc_sys = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.7), Inches(2.5), Inches(1.2))
    kyc_sys.fill.solid()
    kyc_sys.fill.fore_color.rgb = ACCENT_PURPLE
    kyc_sys.line.fill.background()
    tf = kyc_sys.text_frame
    p = tf.paragraphs[0]
    p.text = "KYC시스템\n\n사업자등록증, 등기부\n재무제표"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 화살표 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(1.6), Inches(0.8), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = LIGHT_GRAY
    arrow1.line.fill.background()

    # 화살표 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(3.1), Inches(0.8), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = LIGHT_GRAY
    arrow2.line.fill.background()

    # API 박스들
    api1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(1.4), Inches(2.2), Inches(0.8))
    api1.fill.solid()
    api1.fill.fore_color.rgb = DARK_CARD
    api1.line.color.rgb = ACCENT_BLUE
    api1.line.width = Pt(1)
    tf = api1.text_frame
    p = tf.paragraphs[0]
    p.text = "Snapshot API"
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    api2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(2.9), Inches(2.2), Inches(0.8))
    api2.fill.solid()
    api2.fill.fore_color.rgb = DARK_CARD
    api2.line.color.rgb = ACCENT_PURPLE
    api2.line.width = Pt(1)
    tf = api2.text_frame
    p = tf.paragraphs[0]
    p.text = "Document API"
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_PURPLE
    p.alignment = PP_ALIGN.CENTER

    # 화살표 3, 4
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.6), Inches(1.6), Inches(0.6), Inches(0.4))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = LIGHT_GRAY
    arrow3.line.fill.background()

    arrow4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.6), Inches(3.1), Inches(0.6), Inches(0.4))
    arrow4.fill.solid()
    arrow4.fill.fore_color.rgb = LIGHT_GRAY
    arrow4.line.fill.background()

    # rKYC 플랫폼 (큰 박스)
    rkyc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.4), Inches(1.3), Inches(2.3), Inches(2.8))
    rkyc.fill.solid()
    rkyc.fill.fore_color.rgb = SUCCESS_GREEN
    rkyc.line.fill.background()
    tf = rkyc.text_frame
    p = tf.paragraphs[0]
    p.text = "rKYC\nPlatform\n\n9-Stage\nPipeline"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 출력 화살표
    arrow_out = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(8.35), Inches(4.3), Inches(0.5), Inches(0.6))
    arrow_out.fill.solid()
    arrow_out.fill.fore_color.rgb = SUCCESS_GREEN
    arrow_out.line.fill.background()

    # 출력 박스
    output = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(5.0), Inches(4.2), Inches(1.3))
    output.fill.solid()
    output.fill.fore_color.rgb = DARK_CARD
    output.line.color.rgb = SUCCESS_GREEN
    output.line.width = Pt(2)
    tf = output.text_frame
    p = tf.paragraphs[0]
    p.text = "기업심사 워크플로우\n\n\"이 기업 외부 리스크 시그널 3건 감지\""
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 하단 강조
    highlight = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(5), Inches(0.5))
    tf = highlight.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 기존 인프라 활용 설계가 되어 있음"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    return slide


def add_compliance_slide(prs):
    """규제 준수 방안 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "금융 규제 준수 설계"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 현재 구현된 것
    current_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(4.3), Inches(2.8))
    current_box.fill.solid()
    current_box.fill.fore_color.rgb = DARK_CARD
    current_box.line.color.rgb = SUCCESS_GREEN
    current_box.line.width = Pt(2)

    current_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(4), Inches(0.4))
    tf = current_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ 현재 구현된 것"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN

    current_items = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(4), Inches(2.1))
    tf = current_items.text_frame
    tf.word_wrap = True
    items = [
        "Internal/External LLM 분리 아키텍처",
        "├ 내부 데이터 → Internal LLM만",
        "└ 외부 공개 정보 → External LLM",
        "",
        "LLM Audit Log 테이블",
        "├ 모든 LLM 호출 기록",
        "└ 입력/출력/모델/시간 추적",
        "",
        "Anti-Hallucination 4-Layer",
        "└ 출처 URL 및 감사 추적"
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE if not item.startswith("├") and not item.startswith("└") else LIGHT_GRAY

    # 향후 로드맵
    roadmap_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.1), Inches(4.3), Inches(2.8))
    roadmap_box.fill.solid()
    roadmap_box.fill.fore_color.rgb = DARK_CARD
    roadmap_box.line.color.rgb = ACCENT_BLUE
    roadmap_box.line.width = Pt(2)

    roadmap_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.2), Inches(4), Inches(0.4))
    tf = roadmap_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 향후 로드맵"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    # 로드맵 테이블
    phases = [
        ("MVP (현재)", "외부 API + 추상화", SUCCESS_GREEN),
        ("Phase 2", "Azure/AWS Private", ACCENT_PURPLE),
        ("Phase 3", "On-Premise LLM", DANGER_RED),
    ]

    y_pos = 1.7
    for phase, desc, color in phases:
        phase_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(y_pos), Inches(1.8), Inches(0.6))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = color
        phase_box.line.fill.background()
        tf = phase_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        desc_box = slide.shapes.add_textbox(Inches(7.3), Inches(y_pos), Inches(2.0), Inches(0.6))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        tf.paragraphs[0].space_before = Pt(8)

        y_pos += 0.75

    # 하단 강조
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.3), Inches(7), Inches(0.7))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = ACCENT_BLUE
    highlight.line.fill.background()
    tf = highlight.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 로드맵이 있다는 것 자체가 가점 - 규제 대응 준비됨"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_qa_slide(prs):
    """예상 Q&A 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "심사위원 예상 질문"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Q&A 목록
    qas = [
        ("Q: 외부 LLM에 고객 정보 나가나요?",
         "Internal/External LLM 분리 설계. MVP에서는 공개 정보만 사용,\nPhase 2에서 On-premise LLM 도입 예정"),
        ("Q: 기존 KYC 시스템과 뭐가 다른가요?",
         "기존은 정적 체크리스트. rKYC는 실시간 시그널 모니터링.\n\"지금 이 기업에 무슨 일이 생겼는지\" 알려줌"),
        ("Q: 정확도는?",
         "4-Layer 검증 체계로 Hallucination 방지 설계.\n실제 정확도는 파일럿에서 측정 필요 (솔직히 인정)"),
        ("Q: LLM 죽으면 서비스도 죽나요?",
         "아닙니다. 4-Layer Fallback + Circuit Breaker로\nClaude→GPT→Gemini→Rule-Based 순차 전환"),
    ]

    y_pos = 0.85
    for q, a in qas:
        # 질문
        q_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5))
        q_box.fill.solid()
        q_box.fill.fore_color.rgb = ACCENT_BLUE
        q_box.line.fill.background()
        tf = q_box.text_frame
        p = tf.paragraphs[0]
        p.text = q
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 답변
        a_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.55), Inches(8.8), Inches(0.9))
        tf = a_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "→ " + a
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY

        y_pos += 1.5

    return slide


def add_summary_slide(prs):
    """최종 심사평 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "최종 심사평"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 강점 박스
    strength_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(4.3), Inches(2.5))
    strength_box.fill.solid()
    strength_box.fill.fore_color.rgb = DARK_CARD
    strength_box.line.color.rgb = SUCCESS_GREEN
    strength_box.line.width = Pt(2)

    strength_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(4), Inches(0.4))
    tf = strength_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ 강점"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN

    strengths = [
        "현업 pain point 정확히 이해 (9/10)",
        "해커톤 수준 넘는 완성도 (상위 5%)",
        "확장 가능한 아키텍처",
        "규제 고려한 설계"
    ]
    strength_items = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(4), Inches(1.8))
    tf = strength_items.text_frame
    tf.word_wrap = True
    for i, s in enumerate(strengths):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {s}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)

    # 보완점 박스
    weak_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.1), Inches(4.3), Inches(2.5))
    weak_box.fill.solid()
    weak_box.fill.fore_color.rgb = DARK_CARD
    weak_box.line.color.rgb = WARNING_YELLOW
    weak_box.line.width = Pt(2)

    weak_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.2), Inches(4), Inches(0.4))
    tf = weak_title.text_frame
    p = tf.paragraphs[0]
    p.text = "⚠️ 보완점 (인지)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WARNING_YELLOW

    weaknesses = [
        "실제 정확도 데이터 필요",
        "→ 파일럿에서 측정 예정",
        "",
        "Internal LLM 구현 로드맵",
        "→ Phase 2/3 문서화 완료"
    ]
    weak_items = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(4), Inches(1.8))
    tf = weak_items.text_frame
    tf.word_wrap = True
    for i, w in enumerate(weaknesses):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = w
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE if not w.startswith("→") else LIGHT_GRAY
        p.space_before = Pt(4)

    # 한 줄 요약 (하단 강조)
    summary = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(9), Inches(1.3))
    summary.fill.solid()
    summary.fill.fore_color.rgb = ACCENT_BLUE
    summary.line.fill.background()

    summary_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.15), Inches(8.6), Inches(1.1))
    tf = summary_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 한 줄 요약"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = '"rKYC는 3개 AI Agent가 서로 검증하고, 4단계 Fallback으로\n절대 실패하지 않는, 금융 규제를 충족하는 Enterprise-Grade AI 엔진"'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = WHITE

    return slide

def add_demo_slide(prs):
    """데모 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Live Demo"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # URL 박스
    url_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(8), Inches(1))
    url_box.fill.solid()
    url_box.fill.fore_color.rgb = ACCENT_BLUE
    url_box.line.fill.background()

    tf = url_box.text_frame
    p = tf.paragraphs[0]
    p.text = "https://rkyc-wine.vercel.app/"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 데모 시나리오
    scenarios = [
        "1. Signal Inbox: 전체 시그널 목록 확인",
        "2. 기업 검색: 특정 기업 시그널 필터링",
        "3. Signal Detail: Evidence(출처) 확인",
        "4. Corporate Profile: AI 수집 기업 정보",
        "5. Demo Panel: 분석 실행 트리거"
    ]

    y_pos = 3.0
    for scenario in scenarios:
        box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = scenario
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT_GRAY
        y_pos += 0.55

    return slide

def add_full_architecture_slide(prs):
    """전체 시스템 아키텍처 상세 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Full System Architecture"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ===== 상단: 3-Tier Architecture =====
    # User 아이콘
    user_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(0.9), Inches(0.6), Inches(0.6))
    user_box.fill.solid()
    user_box.fill.fore_color.rgb = LIGHT_GRAY
    user_box.line.fill.background()

    user_label = slide.shapes.add_textbox(Inches(0.15), Inches(1.55), Inches(0.9), Inches(0.3))
    tf = user_label.text_frame
    p = tf.paragraphs[0]
    p.text = "User"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # 화살표 1 (User -> Frontend)
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1.0), Inches(1.05), Inches(0.4), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = LIGHT_GRAY
    arrow1.line.fill.background()

    # Frontend 박스
    frontend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(0.8), Inches(1.8), Inches(0.9))
    frontend.fill.solid()
    frontend.fill.fore_color.rgb = ACCENT_BLUE
    frontend.line.fill.background()
    tf = frontend.text_frame
    p = tf.paragraphs[0]
    p.text = "Frontend\nReact + Vercel"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 화살표 2 (Frontend -> Backend)
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.4), Inches(1.05), Inches(0.4), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = LIGHT_GRAY
    arrow2.line.fill.background()

    # Backend 박스
    backend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.9), Inches(0.8), Inches(1.8), Inches(0.9))
    backend.fill.solid()
    backend.fill.fore_color.rgb = ACCENT_PURPLE
    backend.line.fill.background()
    tf = backend.text_frame
    p = tf.paragraphs[0]
    p.text = "Backend API\nFastAPI + Railway"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 화살표 3 (Backend -> Worker via Redis)
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.8), Inches(1.05), Inches(0.4), Inches(0.3))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = DANGER_RED
    arrow3.line.fill.background()

    # Redis Queue
    redis = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(0.8), Inches(1.2), Inches(0.9))
    redis.fill.solid()
    redis.fill.fore_color.rgb = DANGER_RED
    redis.line.fill.background()
    tf = redis.text_frame
    p = tf.paragraphs[0]
    p.text = "Redis\nQueue"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 화살표 4 (Redis -> Worker)
    arrow4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.6), Inches(1.05), Inches(0.4), Inches(0.3))
    arrow4.fill.solid()
    arrow4.fill.fore_color.rgb = DANGER_RED
    arrow4.line.fill.background()

    # Worker 박스
    worker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(0.8), Inches(1.6), Inches(0.9))
    worker.fill.solid()
    worker.fill.fore_color.rgb = SUCCESS_GREEN
    worker.line.fill.background()
    tf = worker.text_frame
    p = tf.paragraphs[0]
    p.text = "Worker\nCelery"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # ===== 중단: Database & LLM APIs =====
    # Database
    db = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(2.0), Inches(2.0), Inches(0.8))
    db.fill.solid()
    db.fill.fore_color.rgb = WARNING_YELLOW
    db.line.fill.background()
    tf = db.text_frame
    p = tf.paragraphs[0]
    p.text = "Supabase PostgreSQL\n+ pgvector"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER

    # 화살표 (Backend <-> DB)
    arrow_db = slide.shapes.add_shape(MSO_SHAPE.UP_DOWN_ARROW, Inches(4.65), Inches(1.75), Inches(0.25), Inches(0.25))
    arrow_db.fill.solid()
    arrow_db.fill.fore_color.rgb = WARNING_YELLOW
    arrow_db.line.fill.background()

    # ===== 하단: Multi-Agent LLM Layer (큰 박스) =====
    llm_container = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(3.0), Inches(9.4), Inches(3.5))
    llm_container.fill.solid()
    llm_container.fill.fore_color.rgb = RGBColor(25, 33, 52)  # 약간 밝은 배경
    llm_container.line.color.rgb = ACCENT_BLUE
    llm_container.line.width = Pt(2)

    # LLM Layer 제목
    llm_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.4))
    tf = llm_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Agent LLM Layer (Worker Only)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    # Orchestrator 박스 (중앙 상단)
    orchestrator = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3.5), Inches(3), Inches(0.7))
    orchestrator.fill.solid()
    orchestrator.fill.fore_color.rgb = RGBColor(75, 85, 99)
    orchestrator.line.fill.background()
    tf = orchestrator.text_frame
    p = tf.paragraphs[0]
    p.text = "MultiAgentOrchestrator\n(4-Layer Fallback)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 4개 Agent 박스들 (하단 한 줄)
    agents = [
        ("Perplexity\nsonar-pro\n(실시간 검색)", ACCENT_BLUE, 0.5),
        ("Gemini 3 Pro\n(검증/보완)", ACCENT_PURPLE, 2.7),
        ("Claude Opus\n(합성/분석)", SUCCESS_GREEN, 4.9),
        ("OpenAI\nEmbedding\n(벡터화)", WARNING_YELLOW, 7.1),
    ]

    for text, color, x in agents:
        agent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.5), Inches(2.0), Inches(1.0))
        agent.fill.solid()
        agent.fill.fore_color.rgb = color
        agent.line.fill.background()
        tf = agent.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE if color != WARNING_YELLOW else DARK_BG
        p.alignment = PP_ALIGN.CENTER

    # Circuit Breaker & Consensus Engine
    cb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.7), Inches(2.8), Inches(0.6))
    cb.fill.solid()
    cb.fill.fore_color.rgb = DANGER_RED
    cb.line.fill.background()
    tf = cb.text_frame
    p = tf.paragraphs[0]
    p.text = "Circuit Breaker (장애 격리)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    ce = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.7), Inches(3.0), Inches(0.6))
    ce.fill.solid()
    ce.fill.fore_color.rgb = RGBColor(234, 179, 8)  # 더 진한 노랑
    ce.line.fill.background()
    tf = ce.text_frame
    p = tf.paragraphs[0]
    p.text = "Consensus Engine (합의 도출)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER

    cache = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(5.7), Inches(2.8), Inches(0.6))
    cache.fill.solid()
    cache.fill.fore_color.rgb = SUCCESS_GREEN
    cache.line.fill.background()
    tf = cache.text_frame
    p = tf.paragraphs[0]
    p.text = "Profile Cache (7일 TTL)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # ===== 우측 설명 박스 =====
    key_points = slide.shapes.add_textbox(Inches(5.5), Inches(2.0), Inches(4.2), Inches(0.9))
    tf = key_points.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "핵심 설계 원칙"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE

    points = [
        "LLM 키는 Worker에만 존재 (보안)",
        "비동기 처리 (Celery Queue)",
        "모든 AI 호출은 Orchestrator 경유"
    ]
    for point in points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(10)
        p.font.color.rgb = LIGHT_GRAY

    return slide


def add_multi_agent_detail_slide(prs):
    """멀티 에이전트 아키텍처 다이어그램 (계층 구조 스타일)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "rKYC Multi-Agent Architecture"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # ===== Level 1: Entry Point =====
    # run_analysis_pipeline (Celery Task)
    entry_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(0.7), Inches(5), Inches(0.7))
    entry_box.fill.solid()
    entry_box.fill.fore_color.rgb = DARK_CARD
    entry_box.line.color.rgb = LIGHT_GRAY
    entry_box.line.width = Pt(1)
    tf = entry_box.text_frame
    p = tf.paragraphs[0]
    p.text = "run_analysis_pipeline\n(Celery Worker 진입점)"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 수직선 (Entry -> Orchestrator)
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.95), Inches(1.4), Pt(2), Inches(0.3))
    line1.fill.solid()
    line1.fill.fore_color.rgb = LIGHT_GRAY
    line1.line.fill.background()

    # ===== Level 2: Main Orchestrator =====
    orch_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(1.7), Inches(5), Inches(0.7))
    orch_box.fill.solid()
    orch_box.fill.fore_color.rgb = DARK_CARD
    orch_box.line.color.rgb = ACCENT_BLUE
    orch_box.line.width = Pt(2)
    tf = orch_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MultiAgentOrchestrator\n(4-Layer Fallback 조율)"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 분기선 (Orchestrator -> 4 Agents)
    # 중앙 수직선
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.95), Inches(2.4), Pt(2), Inches(0.25))
    line2.fill.solid()
    line2.fill.fore_color.rgb = LIGHT_GRAY
    line2.line.fill.background()

    # 수평선
    line3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.65), Inches(7.6), Pt(2))
    line3.fill.solid()
    line3.fill.fore_color.rgb = LIGHT_GRAY
    line3.line.fill.background()

    # 4개의 수직 분기선
    branch_x = [1.2, 3.4, 5.6, 7.8]
    for x in branch_x:
        branch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.65), Pt(2), Inches(0.2))
        branch.fill.solid()
        branch.fill.fore_color.rgb = LIGHT_GRAY
        branch.line.fill.background()

    # ===== Level 3: 4 Agents =====
    agents = [
        ("PerplexityAgent\n(4 sub-agents)", ACCENT_BLUE, 0.3),
        ("GeminiAgent\n(2 sub-agents)", ACCENT_PURPLE, 2.55),
        ("ClaudeAgent\n(3 sub-agents)", SUCCESS_GREEN, 4.8),
        ("EmbeddingAgent\n(2 sub-agents)", WARNING_YELLOW, 7.05),
    ]

    for text, color, x in agents:
        agent_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.85), Inches(2.1), Inches(0.7))
        agent_box.fill.solid()
        agent_box.fill.fore_color.rgb = DARK_CARD
        agent_box.line.color.rgb = color
        agent_box.line.width = Pt(2)
        tf = agent_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # 수직선 (Agents -> Sub-agents)
    for x in [1.3, 3.55, 5.8, 8.05]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.55), Pt(2), Inches(0.2))
        line.fill.solid()
        line.fill.fore_color.rgb = LIGHT_GRAY
        line.line.fill.background()

    # ===== Level 4: Sub-Agent Details =====
    sub_agents = [
        # Perplexity sub-agents
        ("CorpSearchAgent\nNewsSearchAgent\nIndustrySearchAgent\nPolicySearchAgent", ACCENT_BLUE, 0.3),
        # Gemini sub-agents
        ("ValidatorAgent\nEnricherAgent", ACCENT_PURPLE, 2.55),
        # Claude sub-agents
        ("SynthesizerAgent\nSignalExtractor\nInsightGenerator", SUCCESS_GREEN, 4.8),
        # Embedding sub-agents
        ("TextEmbedder\nSimilaritySearcher", WARNING_YELLOW, 7.05),
    ]

    for text, color, x in sub_agents:
        sub_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.75), Inches(2.1), Inches(1.1))
        sub_box.fill.solid()
        sub_box.fill.fore_color.rgb = DARK_CARD
        sub_box.line.color.rgb = color
        sub_box.line.width = Pt(1)
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

    # 수평선 (Sub-agents -> Pipeline)
    line4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.95), Inches(7.6), Pt(2))
    line4.fill.solid()
    line4.fill.fore_color.rgb = LIGHT_GRAY
    line4.line.fill.background()

    # 수직 분기선들
    for x in [1.3, 3.55, 5.8, 8.05]:
        branch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.85), Pt(2), Inches(0.1))
        branch.fill.solid()
        branch.fill.fore_color.rgb = LIGHT_GRAY
        branch.line.fill.background()

    # 중앙 수직선
    line5 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.95), Inches(4.95), Pt(2), Inches(0.2))
    line5.fill.solid()
    line5.fill.fore_color.rgb = LIGHT_GRAY
    line5.line.fill.background()

    # ===== Level 5: Agent Pipeline =====
    pipeline_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.15), Inches(8.4), Inches(1.2))
    pipeline_box.fill.solid()
    pipeline_box.fill.fore_color.rgb = DARK_CARD
    pipeline_box.line.color.rgb = SUCCESS_GREEN
    pipeline_box.line.width = Pt(2)

    # Pipeline 제목
    pipeline_title = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.4), Inches(0.35))
    tf = pipeline_title.text_frame
    p = tf.paragraphs[0]
    p.text = "9-Stage Analysis Pipeline"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Pipeline 내용
    pipeline_content = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(8.4), Inches(0.8))
    tf = pipeline_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SNAPSHOT → DOC_INGEST → PROFILING → EXTERNAL →"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "CONTEXT → SIGNAL → VALIDATION → INDEX → INSIGHT"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # ===== 우측 하단: Legend =====
    legend_box = slide.shapes.add_textbox(Inches(6.8), Inches(6.45), Inches(3), Inches(0.5))
    tf = legend_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Total: 11 specialized agents"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    return slide


def add_closing_slide(prs):
    """마무리 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Q&A
    qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Q & A"
    p.font.size = Pt(36)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # 로고/태그라인
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "rKYC - Really Know Your Customer"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    """전체 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 타이틀
    add_title_slide(prs, "rKYC", "AI-Powered Corporate Risk Intelligence")

    # 2. 문제 정의 (Pain Point 강화)
    add_content_slide(prs, "Problem: 금융기관 기업 리스크 감지 한계", [
        "외부 정보 수집에 기업당 2시간 소요",
        "뉴스/공시/규제 변화 모니터링이 100% 수작업",
        "KYC 갱신 시점 판단이 담당자 경험에 의존",
        "리스크만 집중하다가 기회 시그널 놓침",
        "",
        "💰 심사역 100명 × 월 10건 × 2시간 = 월 2,000시간 낭비",
        "📊 조기 감지 시 회수율 30-50% 향상 가능"
    ], [5, 6])

    # 3. 솔루션 개요
    add_content_slide(prs, "Solution: rKYC", [
        "Multi-Agent AI 기반 실시간 리스크 감지",
        "",
        "4개 AI 에이전트 협업 (Perplexity + Gemini + Claude + OpenAI)",
        "4-Layer Fallback으로 100% 가용성 보장",
        "Anti-Hallucination 4중 방어로 신뢰성 확보",
        "Production-Grade 엔지니어링 (Circuit Breaker, Redis)"
    ], [2, 3, 4, 5])

    # 4. 전체 시스템 아키텍처
    add_full_architecture_slide(prs)

    # 5. 간단한 아키텍처
    add_architecture_slide(prs)

    # 6. Multi-Agent 개요
    add_multi_agent_slide(prs)

    # 7. Multi-Agent 상세 구현
    add_multi_agent_detail_slide(prs)

    # 8. 4-Layer Fallback
    add_fallback_slide(prs)

    # 9. Anti-Hallucination
    add_anti_hallucination_slide(prs)

    # 10. Pipeline
    add_pipeline_slide(prs)

    # 11. Circuit Breaker
    add_circuit_breaker_slide(prs)

    # 12. 기술 스택
    add_tech_stack_slide(prs)

    # 13. 비즈니스 임팩트 & ROI (강화)
    add_impact_slide(prs)

    # 14. 경쟁 제출작 대비 차별점 (신규)
    add_competitive_slide(prs)

    # 15. 내부 시스템 연계 (신규)
    add_integration_slide(prs)

    # 16. 규제 준수 방안 (신규)
    add_compliance_slide(prs)

    # 17. 예상 Q&A (신규)
    add_qa_slide(prs)

    # 18. 최종 심사평 (신규)
    add_summary_slide(prs)

    # 19. 데모
    add_demo_slide(prs)

    # 20. 마무리
    add_closing_slide(prs)

    return prs


if __name__ == "__main__":
    prs = create_presentation()
    output_path = "rKYC_Hackathon_Pitch.pptx"
    prs.save(output_path)
    print(f"PPT 생성 완료: {output_path}")
