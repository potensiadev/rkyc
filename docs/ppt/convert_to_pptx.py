"""
HTML Slides to PPTX Converter
rKYC AI Engine Architecture Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Color palette
DARK_BG = RgbColor(0x0D, 0x11, 0x17)       # #0D1117
RED_ACCENT = RgbColor(0xE3, 0x37, 0x37)    # #E33737
WHITE = RgbColor(0xFF, 0xFF, 0xFF)          # #FFFFFF
LIGHT_GRAY = RgbColor(0xAA, 0xB7, 0xB8)    # #AAB7B8
DARK_TEXT = RgbColor(0x1C, 0x28, 0x33)     # #1C2833
GREEN = RgbColor(0x27, 0xAE, 0x60)          # #27AE60
BLUE = RgbColor(0x34, 0x98, 0xDB)           # #3498DB
PURPLE = RgbColor(0x9B, 0x59, 0xB6)         # #9B59B6
ORANGE = RgbColor(0xF3, 0x9C, 0x12)         # #F39C12
TEAL = RgbColor(0x16, 0xA0, 0x85)           # #16A085
LIGHT_BG = RgbColor(0xF4, 0xF6, 0xF6)       # #F4F6F6

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_shape_with_text(slide, left, top, width, height, text, bg_color, font_color=WHITE, font_size=14, bold=False, align=PP_ALIGN.CENTER):
    """Add a rounded rectangle with text"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold

    return shape


def add_text_box(slide, left, top, width, height, text, font_color=WHITE, font_size=14, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def create_slide_01_title(prs):
    """Slide 1: Title - rKYC AI Engine Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, DARK_BG)

    # Logo
    add_text_box(slide, Inches(0), Inches(2), SLIDE_WIDTH, Inches(0.6),
                 "rKYC", RED_ACCENT, 24, True, PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, Inches(0), Inches(2.7), SLIDE_WIDTH, Inches(1),
                 "AI Engine Architecture", WHITE, 48, True, PP_ALIGN.CENTER)

    # Engine name
    add_text_box(slide, Inches(0), Inches(3.7), SLIDE_WIDTH, Inches(0.5),
                 "Multi-Agent Orchestration for Risk Intelligence", RED_ACCENT, 22, True, PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(0), Inches(4.3), SLIDE_WIDTH, Inches(0.4),
                 "Enterprise-Grade Corporate Risk Detection Platform", LIGHT_GRAY, 14, False, PP_ALIGN.CENTER)

    # Tech badges
    badges = ["Multi-Agent 4-Layer", "Anti-Hallucination", "9-Stage Pipeline", "Zero-Failure"]
    badge_width = Inches(2.2)
    total_width = len(badges) * badge_width + (len(badges) - 1) * Inches(0.2)
    start_x = (SLIDE_WIDTH - total_width) / 2

    for i, badge in enumerate(badges):
        left = start_x + i * (badge_width + Inches(0.2))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.2), badge_width, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(0x2D, 0x1A, 0x1A)  # Dark red tint
        shape.line.color.rgb = RED_ACCENT
        shape.line.width = Pt(1)

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge
        run.font.size = Pt(11)
        run.font.color.rgb = RED_ACCENT


def create_slide_02_problem(prs):
    """Slide 2: Why AI Engine for KYC?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(6), Inches(0.5),
                 "Why AI Engine for KYC?", WHITE, 28, True)

    # Problem column
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(3), Inches(0.4),
                 "Current Challenges", RED_ACCENT, 18, True)

    problems = [
        ("Manual Risk Monitoring", "1,000+ 기업 수동 모니터링 한계"),
        ("Fragmented Data Sources", "내부 데이터 + 외부 뉴스 + 규제 변화"),
        ("Delayed Signal Detection", "리스크 시그널 조기 탐지 어려움"),
        ("LLM Hallucination Risk", "금융 도메인에서 허위 정보 생성 위험")
    ]

    for i, (title, desc) in enumerate(problems):
        top = Inches(1.8 + i * 1.3)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(5.8), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.fill.background()

        # Left border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(0.05), Inches(1.1))
        border.fill.solid()
        border.fill.fore_color.rgb = RED_ACCENT
        border.line.fill.background()

        add_text_box(slide, Inches(0.7), top + Inches(0.15), Inches(5.5), Inches(0.4),
                     title, DARK_TEXT, 14, True)
        add_text_box(slide, Inches(0.7), top + Inches(0.55), Inches(5.5), Inches(0.4),
                     desc, LIGHT_GRAY, 11, False)

    # Solution column
    add_text_box(slide, Inches(6.8), Inches(1.3), Inches(3), Inches(0.4),
                 "rKYC AI Engine Solution", GREEN, 18, True)

    solutions = [
        ("Multi-Agent Orchestration", "4-Layer Fallback으로 자동화된 분석"),
        ("Unified Context Synthesis", "내부/외부 데이터 통합 컨텍스트"),
        ("Real-time Signal Detection", "9-Stage Pipeline, ~22.7초 분석"),
        ("Anti-Hallucination Defense", "4-Layer 팩트 검증 시스템")
    ]

    for i, (title, desc) in enumerate(solutions):
        top = Inches(1.8 + i * 1.3)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), top, Inches(5.8), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(0xE8, 0xF8, 0xF5)  # Light green
        shape.line.fill.background()

        # Left border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), top, Inches(0.05), Inches(1.1))
        border.fill.solid()
        border.fill.fore_color.rgb = GREEN
        border.line.fill.background()

        add_text_box(slide, Inches(7), top + Inches(0.15), Inches(5.5), Inches(0.4),
                     title, DARK_TEXT, 14, True)
        add_text_box(slide, Inches(7), top + Inches(0.55), Inches(5.5), Inches(0.4),
                     desc, LIGHT_GRAY, 11, False)


def create_slide_03_modules(prs):
    """Slide 3: rKYC AI Engine Core Modules"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
                 "rKYC AI Engine Core Modules", WHITE, 26, True)
    add_text_box(slide, Inches(0.5), Inches(0.7), Inches(8), Inches(0.3),
                 "4 Core Modules for Enterprise Risk Intelligence", RED_ACCENT, 13, False)

    # Modules
    modules = [
        ("Signal Detector", "NLU 기반 리스크 시그널 추출\n3-Tier LLM Fallback Chain", "Claude Opus → GPT-5 → Gemini 3", RED_ACCENT),
        ("Multi-Agent Profiler", "기업 프로파일 자동 구축\n4-Layer Fallback 보장", "Perplexity + Gemini + Consensus", BLUE),
        ("Anti-Hallucination", "4-Layer 팩트 검증\nSource Attribution 필수", "Verify → Guardrail → Validate → Audit", GREEN),
        ("Context Synthesizer", "내부/외부 데이터 통합\n유사 케이스 검색", "pgvector HNSW 2000d", PURPLE)
    ]

    module_width = Inches(2.9)
    start_x = Inches(0.5)

    for i, (name, desc, tech, color) in enumerate(modules):
        left = start_x + i * (module_width + Inches(0.2))

        # Module card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), module_width, Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.fill.background()

        # Top border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), module_width, Inches(0.05))
        border.fill.solid()
        border.fill.fore_color.rgb = color
        border.line.fill.background()

        # Module name
        add_text_box(slide, left + Inches(0.15), Inches(1.6), module_width - Inches(0.3), Inches(0.4),
                     name, color, 15, True)

        # Description
        add_text_box(slide, left + Inches(0.15), Inches(2.1), module_width - Inches(0.3), Inches(1),
                     desc, DARK_TEXT, 11, False)

        # Tech box
        tech_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), Inches(3.5), module_width - Inches(0.3), Inches(0.5))
        tech_box.fill.solid()
        tech_box.fill.fore_color.rgb = WHITE
        tech_box.line.fill.background()

        tf = tech_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = tech
        run.font.size = Pt(9)
        run.font.color.rgb = LIGHT_GRAY

    # Stats
    stats = [
        ("95%+", "Detection Accuracy"),
        ("100%", "Zero-Failure Rate"),
        ("22.7s", "Pipeline Execution"),
        ("7 days", "Profile Cache TTL")
    ]

    stat_width = Inches(2.9)
    for i, (value, label) in enumerate(stats):
        left = start_x + i * (stat_width + Inches(0.2))

        stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.5), stat_width, Inches(1.2))
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = DARK_BG
        stat_box.line.fill.background()

        add_text_box(slide, left, Inches(4.6), stat_width, Inches(0.6),
                     value, RED_ACCENT, 28, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(5.2), stat_width, Inches(0.4),
                     label, LIGHT_GRAY, 12, False, PP_ALIGN.CENTER)


def create_slide_04_architecture(prs):
    """Slide 4: System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
                 "System Architecture", WHITE, 26, True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
                 "3-Layer Architecture + Multi-Agent Pipeline", RED_ACCENT, 12, False)

    # Architecture layers
    layers = [
        ("Presentation Layer", "React 18 + TypeScript + Vite", BLUE),
        ("Business Layer", "FastAPI + SQLAlchemy 2.0", PURPLE),
        ("Data Layer", "PostgreSQL + pgvector", GREEN)
    ]

    layer_width = Inches(3.5)
    start_x = Inches(0.8)

    for i, (name, tech, color) in enumerate(layers):
        left = start_x + i * (layer_width + Inches(0.6))

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), layer_width, Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        add_text_box(slide, left, Inches(1.4), layer_width, Inches(0.4),
                     name, WHITE, 14, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(1.75), layer_width, Inches(0.4),
                     tech, RgbColor(0xFF, 0xFF, 0xFF), 11, False, PP_ALIGN.CENTER)

        if i < 2:
            add_text_box(slide, left + layer_width + Inches(0.1), Inches(1.6), Inches(0.4), Inches(0.4),
                         "→", LIGHT_GRAY, 20, False, PP_ALIGN.CENTER)

    # Arrow down
    add_text_box(slide, Inches(0), Inches(2.4), SLIDE_WIDTH, Inches(0.4),
                 "↓", LIGHT_GRAY, 20, False, PP_ALIGN.CENTER)

    # Worker row
    worker_left = Inches(2)
    worker_width = Inches(4.5)

    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, worker_left, Inches(2.9), worker_width, Inches(0.9))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RED_ACCENT
    box1.line.fill.background()
    add_text_box(slide, worker_left, Inches(3), worker_width, Inches(0.35),
                 "Pipeline Orchestrator", WHITE, 14, True, PP_ALIGN.CENTER)
    add_text_box(slide, worker_left, Inches(3.35), worker_width, Inches(0.35),
                 "Celery 9-Stage + Redis", WHITE, 11, False, PP_ALIGN.CENTER)

    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, worker_left + worker_width + Inches(0.3), Inches(2.9), worker_width, Inches(0.9))
    box2.fill.solid()
    box2.fill.fore_color.rgb = DARK_BG
    box2.line.fill.background()
    add_text_box(slide, worker_left + worker_width + Inches(0.3), Inches(3), worker_width, Inches(0.35),
                 "LLM Ensemble", WHITE, 14, True, PP_ALIGN.CENTER)
    add_text_box(slide, worker_left + worker_width + Inches(0.3), Inches(3.35), worker_width, Inches(0.35),
                 "Claude Opus 4.5 | GPT-5 | Gemini 3 Pro", WHITE, 11, False, PP_ALIGN.CENTER)

    # Arrow down
    add_text_box(slide, Inches(0), Inches(3.9), SLIDE_WIDTH, Inches(0.4),
                 "↓", LIGHT_GRAY, 20, False, PP_ALIGN.CENTER)

    # External services
    services = [
        ("Web Intelligence", "Perplexity sonar-pro", ORANGE),
        ("Embedding Engine", "text-embedding-3-large", TEAL),
        ("Vector Store", "pgvector HNSW 2000d", RgbColor(0x8E, 0x44, 0xAD))
    ]

    svc_width = Inches(3.5)
    for i, (name, tech, color) in enumerate(services):
        left = start_x + i * (svc_width + Inches(0.4))

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.4), svc_width, Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        add_text_box(slide, left, Inches(4.5), svc_width, Inches(0.35),
                     name, WHITE, 12, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(4.8), svc_width, Inches(0.3),
                     tech, RgbColor(0xFF, 0xFF, 0xFF), 10, False, PP_ALIGN.CENTER)

    # Deploy info
    deploys = [
        ("Frontend: Vercel", "rkyc-wine.vercel.app"),
        ("Backend: Railway", "rkyc-production.up.railway.app"),
        ("Database: Supabase", "Tokyo (ap-northeast-1)"),
        ("Worker: Railway", "Celery + Redis")
    ]

    deploy_width = Inches(2.9)
    for i, (name, url) in enumerate(deploys):
        left = Inches(0.5) + i * (deploy_width + Inches(0.2))

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.5), deploy_width, Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.fill.background()

        add_text_box(slide, left, Inches(5.55), deploy_width, Inches(0.3),
                     name, DARK_TEXT, 10, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(5.85), deploy_width, Inches(0.3),
                     url, LIGHT_GRAY, 8, False, PP_ALIGN.CENTER)


def create_slide_05_ensemble(prs):
    """Slide 5: Multi-Model Ensemble"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
                 "Multi-Model Ensemble", WHITE, 26, True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
                 "3-Tier LLM Fallback Chain for High Availability", RED_ACCENT, 12, False)

    # Labels
    labels = ["PRIMARY", "FALLBACK 1", "FALLBACK 2"]
    model_width = Inches(3.5)
    start_x = Inches(0.8)

    for i, label in enumerate(labels):
        left = start_x + i * (model_width + Inches(0.6))
        add_text_box(slide, left, Inches(1.3), model_width, Inches(0.3),
                     label, LIGHT_GRAY, 11, False, PP_ALIGN.CENTER)

    # Model boxes
    models = [
        ("Claude Opus 4.5", "claude-opus-4-5-20251101", "200K context | Complex reasoning", RED_ACCENT),
        ("GPT-5", "gpt-5", "128K context | Structured output", BLUE),
        ("Gemini 3 Pro", "gemini-3-pro-preview", "1M context | Cross-validation", GREEN)
    ]

    for i, (name, model_id, spec, color) in enumerate(models):
        left = start_x + i * (model_width + Inches(0.6))

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), model_width, Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        add_text_box(slide, left, Inches(1.85), model_width, Inches(0.4),
                     name, WHITE, 16, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.25), model_width, Inches(0.3),
                     model_id, RgbColor(0xFF, 0xFF, 0xFF), 10, False, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.7), model_width, Inches(0.3),
                     spec, RgbColor(0xFF, 0xFF, 0xFF), 9, False, PP_ALIGN.CENTER)

        if i < 2:
            add_text_box(slide, left + model_width + Inches(0.1), Inches(2.2), Inches(0.4), Inches(0.4),
                         "→", LIGHT_GRAY, 24, True, PP_ALIGN.CENTER)

    # Fallback info
    info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.5), Inches(12.3), Inches(2))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = LIGHT_BG
    info_box.line.fill.background()

    add_text_box(slide, Inches(0.7), Inches(3.65), Inches(4), Inches(0.4),
                 "Fallback Trigger Conditions", DARK_BG, 14, True)

    conditions = [
        ("API Timeout", "30초 초과 시 다음 모델로 전환"),
        ("Rate Limit", "429 에러 발생 시 자동 전환"),
        ("Invalid Response", "JSON 파싱 실패 시 재시도"),
        ("Circuit Breaker", "연속 3회 실패 시 5분 cooldown")
    ]

    cond_width = Inches(2.8)
    for i, (title, desc) in enumerate(conditions):
        left = Inches(0.7) + i * (cond_width + Inches(0.2))

        cond_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(4.2), cond_width, Inches(1))
        cond_box.fill.solid()
        cond_box.fill.fore_color.rgb = WHITE
        cond_box.line.fill.background()

        # Left border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(4.2), Inches(0.05), Inches(1))
        border.fill.solid()
        border.fill.fore_color.rgb = RED_ACCENT
        border.line.fill.background()

        add_text_box(slide, left + Inches(0.15), Inches(4.3), cond_width - Inches(0.2), Inches(0.3),
                     title, DARK_TEXT, 11, True)
        add_text_box(slide, left + Inches(0.15), Inches(4.65), cond_width - Inches(0.2), Inches(0.4),
                     desc, LIGHT_GRAY, 9, False)


def create_slide_06_orchestrator(prs):
    """Slide 6: Multi-Agent Orchestrator"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
                 "Multi-Agent Orchestrator", WHITE, 26, True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
                 "4-Layer Fallback System for Zero-Failure Guarantee", WHITE, 12, False)

    # Layers
    layers = [
        [("L0", "Intelligent Cache", "7-day TTL, Profile Store", RgbColor(0xF4, 0xF6, 0xF6)),
         ("L1", "Perplexity Search", "sonar-pro, Citation with URLs", RgbColor(0xE8, 0xF8, 0xF5))],
        [("L1.5", "Gemini Validation", "Cross-check, Enrich missing fields", RgbColor(0xFE, 0xF9, 0xE7)),
         ("L2", "Consensus Engine", "Jaccard Similarity >= 0.7", RgbColor(0xEB, 0xF5, 0xFB))],
        [("L3", "Rule-Based Merge", "Source priority, Range validation", RgbColor(0xF5, 0xEE, 0xF8)),
         ("L4", "Graceful Degradation", "Minimal profile, _degraded flag", RgbColor(0xFD, 0xED, 0xEC))]
    ]

    layer_width = Inches(5.8)

    for row_idx, row in enumerate(layers):
        for col_idx, (layer_id, title, desc, color) in enumerate(row):
            left = Inches(0.5) + col_idx * (layer_width + Inches(0.3))
            top = Inches(1.3) + row_idx * Inches(1.4)

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, layer_width, Inches(1))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            # Layer number circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.2), Inches(0.6), Inches(0.6))
            circle.fill.solid()
            circle.fill.fore_color.rgb = DARK_BG
            circle.line.fill.background()

            add_text_box(slide, left + Inches(0.15), top + Inches(0.35), Inches(0.6), Inches(0.3),
                         layer_id, WHITE, 12, True, PP_ALIGN.CENTER)

            add_text_box(slide, left + Inches(0.9), top + Inches(0.2), layer_width - Inches(1.1), Inches(0.35),
                         title, DARK_TEXT, 14, True)
            add_text_box(slide, left + Inches(0.9), top + Inches(0.55), layer_width - Inches(1.1), Inches(0.35),
                         desc, DARK_TEXT, 10, False)

        if row_idx < 2:
            add_text_box(slide, Inches(0), top + Inches(1.05), SLIDE_WIDTH, Inches(0.3),
                         "↓", LIGHT_GRAY, 16, False, PP_ALIGN.CENTER)

    # Guarantee stats
    stats = [
        ("100%", "Zero-Failure Rate"),
        ("7 days", "Cache TTL"),
        ("1 day", "Fallback TTL"),
        ("19", "Profile Fields")
    ]

    stat_width = Inches(2.5)
    start_x = Inches(1.2)

    for i, (value, label) in enumerate(stats):
        left = start_x + i * (stat_width + Inches(0.3))

        stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.6), stat_width, Inches(1))
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = DARK_BG
        stat_box.line.fill.background()

        add_text_box(slide, left, Inches(5.7), stat_width, Inches(0.5),
                     value, RED_ACCENT, 22, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(6.15), stat_width, Inches(0.3),
                     label, LIGHT_GRAY, 10, False, PP_ALIGN.CENTER)


def create_slide_07_antihallucination(prs):
    """Slide 7: Anti-Hallucination Defense"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = RED_ACCENT
    header.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
                 "Anti-Hallucination Defense", WHITE, 26, True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
                 "4-Layer Fact Verification System", WHITE, 12, False)

    # Defense layers
    layers = [
        ("L1", "Source Verification", "Domain Credibility Scorer (T1: 공시/IR → T4: 블로그)", RgbColor(0xE8, 0xF8, 0xF5)),
        ("L2", "Extraction Guardrails", '"Return null if uncertain" Policy, Schema Compliance', RgbColor(0xFE, 0xF9, 0xE7)),
        ("L3", "Cross-Validation", "Gemini Validator, Jaccard Similarity >= 0.7, Discrepancy Flag", RgbColor(0xEB, 0xF5, 0xFB)),
        ("L4", "Audit Trail", "field_provenance JSON, LLM Audit Log, Trace ID", RgbColor(0xF5, 0xEE, 0xF8))
    ]

    layer_width = Inches(7.5)

    for i, (layer_id, title, desc, color) in enumerate(layers):
        top = Inches(1.3) + i * Inches(1.2)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, layer_width, Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Layer number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.2), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = DARK_BG
        circle.line.fill.background()

        add_text_box(slide, Inches(0.7), top + Inches(0.35), Inches(0.6), Inches(0.3),
                     layer_id, WHITE, 14, True, PP_ALIGN.CENTER)

        add_text_box(slide, Inches(1.5), top + Inches(0.2), Inches(6), Inches(0.35),
                     title, DARK_TEXT, 14, True)
        add_text_box(slide, Inches(1.5), top + Inches(0.55), Inches(6), Inches(0.35),
                     desc, DARK_TEXT, 10, False)

    # Provenance section
    prov_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.3), Inches(4.5), Inches(5.3))
    prov_box.fill.solid()
    prov_box.fill.fore_color.rgb = DARK_BG
    prov_box.line.fill.background()

    add_text_box(slide, Inches(8.5), Inches(1.5), Inches(4), Inches(0.4),
                 "Provenance Tracking", RED_ACCENT, 14, True)

    prov_items = [
        ("source_url", "출처 URL (필수)"),
        ("excerpt", "원문 발췌 (50자 이내)"),
        ("confidence", "HIGH / MED / LOW"),
        ("retrieved_at", "검색 시점 타임스탬프")
    ]

    for i, (key, desc) in enumerate(prov_items):
        top = Inches(2) + i * Inches(0.8)

        item_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), top, Inches(4.1), Inches(0.65))
        item_box.fill.solid()
        item_box.fill.fore_color.rgb = RgbColor(0x1C, 0x28, 0x33)
        item_box.line.fill.background()

        add_text_box(slide, Inches(8.6), top + Inches(0.1), Inches(3.9), Inches(0.3),
                     key, WHITE, 11, True)
        add_text_box(slide, Inches(8.6), top + Inches(0.35), Inches(3.9), Inches(0.25),
                     desc, LIGHT_GRAY, 9, False)

    # Confidence levels
    add_text_box(slide, Inches(8.5), Inches(5.3), Inches(4), Inches(0.3),
                 "Confidence Level", ORANGE, 12, True)

    conf_colors = [GREEN, ORANGE, RED_ACCENT]
    conf_labels = ["HIGH", "MED", "LOW"]

    for i, (color, label) in enumerate(zip(conf_colors, conf_labels)):
        left = Inches(8.5) + i * Inches(1.3)

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.7), Inches(1.1), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()

        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(10)
        run.font.color.rgb = WHITE
        run.font.bold = True

    add_text_box(slide, Inches(8.5), Inches(6.2), Inches(4), Inches(0.3),
                 "공시/IR > 뉴스 > 추정", LIGHT_GRAY, 9, False)


def create_slide_08_pipeline(prs):
    """Slide 8: 9-Stage Analysis Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BG
    header.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
                 "9-Stage Analysis Pipeline", WHITE, 26, True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
                 "Automated Risk Signal Detection | ~22.7s Complete", RED_ACCENT, 12, False)

    # Pipeline stages
    stages = [
        ("1", "SNAPSHOT", "Internal Load", False),
        ("2", "DOC", "PDF Parser", False),
        ("3 ★", "PROFILE", "Multi-Agent", True),
        ("4", "EXTERNAL", "Perplexity", False),
        ("5", "CONTEXT", "Synthesis", False),
        ("6", "SIGNAL", "LLM Extract", False),
        ("7 ★", "VALIDATE", "Anti-Hallu", True),
        ("8", "INDEX", "Vector 2000d", False),
        ("9", "INSIGHT", "Briefing", False)
    ]

    stage_width = Inches(1.2)
    start_x = Inches(0.3)

    for i, (num, name, desc, is_key) in enumerate(stages):
        left = start_x + i * (stage_width + Inches(0.15))

        color = RgbColor(0xE8, 0xF8, 0xF5) if is_key else LIGHT_BG
        border_color = GREEN if is_key else RED_ACCENT

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), stage_width, Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Left border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.3), Inches(0.03), Inches(1.4))
        border.fill.solid()
        border.fill.fore_color.rgb = border_color
        border.line.fill.background()

        num_color = GREEN if is_key else RED_ACCENT
        add_text_box(slide, left, Inches(1.4), stage_width, Inches(0.25),
                     num, num_color, 10, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(1.7), stage_width, Inches(0.3),
                     name, DARK_TEXT, 10, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.05), stage_width, Inches(0.3),
                     desc, LIGHT_GRAY, 8, False, PP_ALIGN.CENTER)

        if i < 8:
            add_text_box(slide, left + stage_width, Inches(1.8), Inches(0.15), Inches(0.3),
                         "→", LIGHT_GRAY, 12, False, PP_ALIGN.CENTER)

    # Bottom info boxes
    info_sections = [
        ("LLM Ensemble", ["Claude Opus 4.5 (Primary)", "GPT-5 (Fallback 1)", "Gemini 3 Pro (Fallback 2)"]),
        ("Resilience", ["Circuit Breaker (3 fail / 5min)", "Rate Limiter (10/min)", "Auto Recovery"]),
        ("Vector Search", ["pgvector HNSW 2000d", "Cosine Similarity", "Case-Based Reasoning"])
    ]

    box_width = Inches(4)
    for i, (title, items) in enumerate(info_sections):
        left = Inches(0.5) + i * (box_width + Inches(0.2))

        info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3), box_width, Inches(2.2))
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = DARK_BG
        info_box.line.fill.background()

        add_text_box(slide, left + Inches(0.15), Inches(3.15), box_width - Inches(0.3), Inches(0.35),
                     title, RED_ACCENT, 12, True)

        for j, item in enumerate(items):
            item_top = Inches(3.55) + j * Inches(0.5)

            item_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), item_top, box_width - Inches(0.3), Inches(0.4))
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = RgbColor(0x1C, 0x28, 0x33)
            item_box.line.fill.background()

            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item
            run.font.size = Pt(9)
            run.font.color.rgb = WHITE


def create_slide_09_database(prs):
    """Slide 9: Data & Vector Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = GREEN
    header.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.25), Inches(8), Inches(0.5),
                 "Data & Vector Architecture", WHITE, 26, True)
    add_text_box(slide, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
                 "PostgreSQL + pgvector | Semantic Search Engine", WHITE, 12, False)

    # Table groups
    groups = [
        ("Core Master", ["corp", "industry_master"]),
        ("Signal System", ["rkyc_signal", "rkyc_signal_index", "rkyc_evidence", "rkyc_signal_embedding"]),
        ("Profile & Context", ["rkyc_corp_profile", "rkyc_internal_snapshot", "rkyc_unified_context"]),
        ("Document & External", ["rkyc_document", "rkyc_fact", "rkyc_external_news"]),
        ("Audit & Jobs", ["rkyc_llm_audit_log", "rkyc_job", "rkyc_case_index"])
    ]

    left = Inches(0.5)
    top = Inches(1.2)

    for group_name, tables in groups:
        add_text_box(slide, left, top, Inches(3), Inches(0.3),
                     group_name, RED_ACCENT, 12, True)
        top += Inches(0.35)

        table_left = left
        for table in tables:
            table_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, table_left, top, Inches(1.8), Inches(0.4))
            table_box.fill.solid()
            table_box.fill.fore_color.rgb = LIGHT_BG
            table_box.line.fill.background()

            tf = table_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = table
            run.font.size = Pt(9)
            run.font.color.rgb = DARK_TEXT

            table_left += Inches(1.9)
            if table_left > Inches(7):
                table_left = left
                top += Inches(0.45)

        top += Inches(0.6)

    # Vector Store section
    vec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.2), Inches(4.5), Inches(5.3))
    vec_box.fill.solid()
    vec_box.fill.fore_color.rgb = DARK_BG
    vec_box.line.fill.background()

    add_text_box(slide, Inches(8.5), Inches(1.4), Inches(4), Inches(0.4),
                 "Vector Store", GREEN, 14, True)

    vec_items = [
        ("Embedding Model", "text-embedding-3-large (2000d)"),
        ("Index Type", "HNSW (m=16, ef_construction=64)"),
        ("Similarity", "Cosine Distance (vector <=> vector)")
    ]

    for i, (title, desc) in enumerate(vec_items):
        top = Inches(1.9) + i * Inches(1)

        item_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), top, Inches(4.1), Inches(0.8))
        item_box.fill.solid()
        item_box.fill.fore_color.rgb = RgbColor(0x1C, 0x28, 0x33)
        item_box.line.fill.background()

        add_text_box(slide, Inches(8.6), top + Inches(0.1), Inches(3.9), Inches(0.3),
                     title, WHITE, 12, True)
        add_text_box(slide, Inches(8.6), top + Inches(0.4), Inches(3.9), Inches(0.3),
                     desc, LIGHT_GRAY, 10, False)

    # Key features
    add_text_box(slide, Inches(8.5), Inches(5), Inches(4), Inches(0.35),
                 "Key Features", ORANGE, 12, True)

    features = [
        "Signal Embedding (rkyc_signal_embedding)",
        "Case-Based Reasoning (rkyc_case_index)",
        "Semantic Search for Similar Cases",
        "Provenance JSON Pointer"
    ]

    for i, feature in enumerate(features):
        add_text_box(slide, Inches(8.5), Inches(5.4) + i * Inches(0.35), Inches(4), Inches(0.3),
                     feature, LIGHT_GRAY, 9, False)


def create_slide_10_demo(prs):
    """Slide 10: Live Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Accent bar at bottom
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RED_ACCENT
    accent_bar.line.fill.background()

    # Title
    add_text_box(slide, Inches(0), Inches(2), SLIDE_WIDTH, Inches(0.8),
                 "Live Demo", WHITE, 44, True, PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(0.5),
                 "rKYC AI Engine in Action", RED_ACCENT, 20, True, PP_ALIGN.CENTER)

    # Engine badge
    add_text_box(slide, Inches(0), Inches(3.3), SLIDE_WIDTH, Inches(0.4),
                 "Multi-Agent Orchestration | Anti-Hallucination | 9-Stage Pipeline", LIGHT_GRAY, 12, False, PP_ALIGN.CENTER)

    # URL boxes
    urls = [
        ("Platform", "rkyc-wine.vercel.app"),
        ("API Endpoint", "rkyc-production.up.railway.app")
    ]

    url_width = Inches(4)
    start_x = (SLIDE_WIDTH - 2 * url_width - Inches(0.5)) / 2

    for i, (label, url) in enumerate(urls):
        left = start_x + i * (url_width + Inches(0.5))

        url_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4), url_width, Inches(0.9))
        url_box.fill.solid()
        url_box.fill.fore_color.rgb = RgbColor(0x1C, 0x28, 0x33)
        url_box.line.fill.background()

        add_text_box(slide, left, Inches(4.1), url_width, Inches(0.3),
                     label, LIGHT_GRAY, 11, False, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(4.4), url_width, Inches(0.35),
                     url, WHITE, 13, False, PP_ALIGN.CENTER)

    # Feature stats
    features = [
        ("6", "Corporations"),
        ("29", "Signals"),
        ("22.7s", "Pipeline"),
        ("95%+", "Accuracy"),
        ("100%", "Zero-Failure")
    ]

    feat_width = Inches(2)
    total_width = len(features) * feat_width + (len(features) - 1) * Inches(0.3)
    start_x = (SLIDE_WIDTH - total_width) / 2

    for i, (value, label) in enumerate(features):
        left = start_x + i * (feat_width + Inches(0.3))

        add_text_box(slide, left, Inches(5.3), feat_width, Inches(0.6),
                     value, RED_ACCENT, 26, True, PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(5.9), feat_width, Inches(0.3),
                     label, LIGHT_GRAY, 11, False, PP_ALIGN.CENTER)


def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Create slides
    print("Creating slides...")
    create_slide_01_title(prs)
    print("  - Slide 1: Title")
    create_slide_02_problem(prs)
    print("  - Slide 2: Why AI Engine")
    create_slide_03_modules(prs)
    print("  - Slide 3: Core Modules")
    create_slide_04_architecture(prs)
    print("  - Slide 4: System Architecture")
    create_slide_05_ensemble(prs)
    print("  - Slide 5: Multi-Model Ensemble")
    create_slide_06_orchestrator(prs)
    print("  - Slide 6: Multi-Agent Orchestrator")
    create_slide_07_antihallucination(prs)
    print("  - Slide 7: Anti-Hallucination Defense")
    create_slide_08_pipeline(prs)
    print("  - Slide 8: 9-Stage Pipeline")
    create_slide_09_database(prs)
    print("  - Slide 9: Data & Vector Architecture")
    create_slide_10_demo(prs)
    print("  - Slide 10: Live Demo")

    # Save presentation
    output_path = "rKYC_AI_Engine_Architecture.pptx"
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")


if __name__ == "__main__":
    main()
